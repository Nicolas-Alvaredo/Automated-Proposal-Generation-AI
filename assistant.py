import os
import time
import requests
import json
from openai import AzureOpenAI
from dotenv import load_dotenv
from docx import Document as DocxDocument
import pandas as pd
import pdfplumber
import pptx
import msal
import io
import logging
import jwt
from azure.identity import DefaultAzureCredential
from azure.storage.blob import BlobServiceClient
import azure.functions as func
import tempfile
from win32com import client as win32  # Import win32com

# Adjust the logging level to reduce verbosity
logging.basicConfig(level=logging.INFO)
httpx_log = logging.getLogger("httpx")
httpx_log.setLevel(logging.WARNING)

def get_access_token(client_id, client_secret, tenant_id, scope):
    authority = f"https://login.microsoftonline.com/{tenant_id}"
    app = msal.ConfidentialClientApplication(client_id, authority=authority, client_credential=client_secret)
    result = app.acquire_token_for_client(scopes=[scope])

    if "access_token" not in result:
        raise Exception(f"Could not obtain access token for scope: {scope}")

    access_token = result["access_token"]
    logging.info(f"Obtained access token for scope {scope}: {access_token}")
    
    # Decode and log token info
    token_info = jwt.decode(access_token, options={"verify_signature": False})
    logging.info(f"Token info for scope {scope}: {token_info}")
    
    return access_token

def create_openai_client(api_key, api_version, azure_endpoint):
    return AzureOpenAI(
        api_key=api_key,
        api_version=api_version,
        azure_endpoint=azure_endpoint
    )

def handle_request(attachments_folder_id, response_folder_id):
    try:
        # Load environment variables from .env file
        load_dotenv()

        logging.info(f"Attachments Folder ID: {attachments_folder_id}")
        logging.info(f"Response Folder ID: {response_folder_id}")

        # Load environment variables for OpenAI and Microsoft Graph
        api_key = os.getenv("AZURE_OPENAI_API_KEY")
        api_version = "2024-02-15-preview"
        azure_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
        client_id = os.getenv("MS_GRAPH_CLIENT_ID")
        client_secret = os.getenv("MS_GRAPH_CLIENT_SECRET")
        tenant_id = os.getenv("MS_GRAPH_TENANT_ID")
        drive_id = os.getenv("MS_GRAPH_DRIVE_ID")
        adls_account_name = os.getenv("ADLS_ACCOUNT_NAME")
        adls_container_name = os.getenv("ADLS_CONTAINER_NAME")

        # Validate environment variables
        if not api_key or not azure_endpoint:
            raise ValueError("Missing required environment variables for OpenAI. Please check your .env file.")
        if not client_id or not client_secret or not tenant_id or not drive_id:
            raise ValueError("Missing required environment variables for Microsoft Graph. Please check your .env file.")
        if not adls_account_name or not adls_container_name:
            raise ValueError("Missing required environment variables for ADLS. Please check your .env file.")

        # Get access token for Microsoft Graph
        graph_access_token = get_access_token(client_id, client_secret, tenant_id, "https://graph.microsoft.com/.default")

        # Initialize the Azure OpenAI client
        client = create_openai_client(api_key, api_version, azure_endpoint)

        logging.info("Listing files in the Attachments folder")
        headers = {'Authorization': f'Bearer {graph_access_token}'}
        list_files_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{attachments_folder_id}/children"
        response = requests.get(list_files_url, headers=headers)

        if response.status_code != 200:
            raise Exception(f"Failed to list files in the Attachments folder: {response.status_code} {response.text}")

        files = response.json().get('value', [])
        file_contents = []

        for file in files:
            file_name = file['name']
            file_id = file['id']
            download_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{file_id}/content"
            file_response = requests.get(download_url, headers=headers)
            if file_response.status_code == 200:
                file_content = file_response.content
                if file_name.lower().endswith('.doc'):
                    file_name = file_name[:-4] + '.docx'
                    file_content = convert_doc_to_docx(file_content)
                upload_to_adls(file_name, file_content, attachments_folder_id, adls_account_name, adls_container_name)
            else:
                logging.error(f"Failed to download file '{file_name}': {file_response.status_code}")

        logging.info("Downloading files from ADLS")
        adls_file_contents = download_from_adls(attachments_folder_id, adls_account_name, adls_container_name)
        
        logging.info("Creating chunks for each file")
        chunk_size = 100000  # Adjust the chunk size to be within the limit
        chunks = []
        for file_content in adls_file_contents:
            chunks.extend([file_content[i:i + chunk_size] for i in range(0, len(file_content), chunk_size)])

        logging.info("Preparing prompts")
        prompts = [f"Please acknowledge and save the following content chunk {i + 1}:\n{chunk}" for i, chunk in enumerate(chunks)]
        prompts.append("RFP Analysis Prompt: Please thoroughly analyze the attached RFP document, including all amendments, Q&A responses, and related attachments. Identify and extract the following key information:\n\nCustomer's primary objectives and requirements\n\nEvaluation criteria and their relative weightings\n\nSubmission deadlines and formatting requirements\n\nAny unique or mandatory compliance requirements\n\nSummarize your findings in a clear, concise report highlighting the most critical elements we must address to be fully responsive and compliant.")
        prompts.append("Customer Problem/Objective Identification Prompt: Based on your analysis of the RFP and any additional customer background information provided, identify the customer's top 3-5 problems, pain points, or strategic objectives that our solution must address to win. For each problem/objective:\n\nDescribe the current situation and its negative impacts on the customer\n\nHighlight the urgency and importance of addressing it\n\nIdentify any specific metrics or success criteria the customer has defined\n\nSuggest potential solution elements or approaches that could effectively tackle the problem\n\nPresent your findings in a prioritized list with supporting rationale for each item.")
        prompts.append("Full Proposal Draft Assembly Prompt: Please assemble the complete first draft of the proposal, integrating all AI-generated section content according to the approved outline structure. Ensure that:\n\nAll sections and subsections flow logically and persuasively, with clear transitions and cross-references as needed\n\nAll RFP requirements and evaluation criteria are fully addressed, with no gaps or redundancies\n\nAll win themes, differentiators, and proof points are consistently messaged and mutually reinforcing across sections\n\nContinuity and consistency across all sections in terms of customer focus, tone, style, and reading level\n\nPlaceholders for graphics, tables, and callout boxes are appropriate and properly formatted\n\nAll required attachments, forms, and administrative elements are included and compliant\n\nPlease provide a detailed table of contents and cross-reference matrix to aid in navigation and compliance reviews. Clearly label any areas requiring further SME input or validation.")
        prompts.append("Please, according to all consolidated info from previous prompts, create the final proposal (With right format and titles, please write titles and subtitles between **)")

        logging.info("Creating thread and adding prompts")
        try:
            thread = client.beta.threads.create()
            thread_id = thread.id
            logging.info(f"Thread created with ID: {thread_id}")
        except Exception as e:
            logging.error(f"Failed to create thread: {e}")
            raise

        chat_doc = DocxDocument()
        chat_doc.add_heading('Chat History', 0)

        for i, prompt in enumerate(prompts):
            attempt = 0
            success = False
            while attempt < 3 and not success:
                try:
                    message = client.beta.threads.messages.create(
                        thread_id=thread_id,
                        role="user",
                        content=prompt
                    )
                    logging.info(f"Added prompt {i+1}/{len(prompts)} to thread")
                    success = True
                except Exception as e:
                    logging.error(f"Failed to add prompt {i+1}/{len(prompts)} to thread: {e}")
                    attempt += 1
                    time.sleep(1)  # Delay before retrying

            if not success:
                logging.error(f"Giving up on prompt {i+1}/{len(prompts)}: {prompt}")
                continue

            attempt = 0
            success = False
            while attempt < 3 and not success:
                try:
                    run = client.beta.threads.runs.create(
                        thread_id=thread_id,
                        assistant_id="asst_Wcvq39x7iEhLkyTOxGuoMQR7"
                    )
                    success = True
                except Exception as e:
                    logging.error(f"Failed to run thread for prompt {i+1}/{len(prompts)}: {e}")
                    attempt += 1
                    time.sleep(1)  # Delay before retrying

            if not success:
                logging.error(f"Giving up on running thread for prompt {i+1}/{len(prompts)}")
                continue

            while run.status in ['queued', 'in_progress', 'cancelling']:
                time.sleep(1)
                run = client.beta.threads.runs.retrieve(
                    thread_id=thread_id,
                    run_id=run.id
                )

            if run.status != 'completed':
                logging.error(f"Run for question {i+1}/{len(prompts)} failed with status: {run.status}")
                continue

            # Retrieve the message and add it to the chat history
            messages = client.beta.threads.messages.list(
                thread_id=thread_id
            )

            sorted_messages = sorted(messages.data, key=lambda x: x.created_at)

            for message in sorted_messages:
                role = "User" if message.role == "user" else "Assistant"
                content = message.content[0].text.value

                p = chat_doc.add_paragraph()
                p.add_run(f"{role}: ").bold = True
                p.add_run(content)

        chat_byte_stream = io.BytesIO()
        chat_doc.save(chat_byte_stream)
        chat_byte_stream.seek(0)

        draft_doc = DocxDocument()
        draft_doc.add_heading('Draft Proposal', 0)
        last_response = sorted_messages[-1].content[0].text.value if sorted_messages else "No response available."

        # Add formatted content to the draft document
        add_formatted_content(draft_doc, last_response)

        draft_byte_stream = io.BytesIO()
        draft_doc.save(draft_byte_stream)
        draft_byte_stream.seek(0)

        upload_to_onedrive(chat_byte_stream, 'chat_history.docx', response_folder_id, drive_id, graph_access_token)
        upload_to_onedrive(draft_byte_stream, 'draft_response.docx', response_folder_id, drive_id, graph_access_token)
    except Exception as e:
        logging.error(f"Error in handle_request: {e}", exc_info=True)
        raise

def process_file(file_name, file_content):
    file_ext = os.path.splitext(file_name)[1].lower()
    if file_ext == '.xlsx' or file_ext == '.xls':
        return process_excel(file_content)
    elif file_ext == '.docx' or file_ext == '.doc':
        if file_ext == '.doc':
            file_content = convert_doc_to_docx(file_content)
        return process_word(file_content)
    elif file_ext == '.pdf':
        return process_pdf(file_content)
    elif file_ext == '.pptx':
        return process_ppt(file_content)
    elif file_ext == '.txt':
        return process_text(file_content)
    else:
        return f"Unsupported file type: {file_ext}"

def process_excel(file_content):
    try:
        df = pd.read_excel(io.BytesIO(file_content))
        return df.to_string()
    except Exception as e:
        logging.error(f"Error processing Excel file: {e}")
        return f"Error processing Excel file: {e}"

def process_word(file_content):
    doc = DocxDocument(io.BytesIO(file_content))
    paragraphs = [p.text for p in doc.paragraphs]
    return "\n".join(paragraphs)

def process_pdf(file_content):
    text = ""
    with pdfplumber.open(io.BytesIO(file_content)) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

def process_ppt(file_content):
    prs = pptx.Presentation(io.BytesIO(file_content))
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def process_text(file_content):
    return file_content.decode('utf-8')

def add_formatted_content(draft_doc, content):
    # Split the content into paragraphs and add formatting
    paragraphs = content.split('\n')
    for paragraph in paragraphs:
        p = draft_doc.add_paragraph()
        while '**' in paragraph:
            start_index = paragraph.find('**')
            end_index = paragraph.find('**', start_index + 2)
            if end_index == -1:
                break
            # Add text before the bold part
            p.add_run(paragraph[:start_index])
            # Add bold text
            p.add_run(paragraph[start_index + 2:end_index]).bold = True
            # Update paragraph to the remaining part
            paragraph = paragraph[end_index + 2:]
        # Add any remaining text after the last bold part
        p.add_run(paragraph)

def upload_to_adls(file_name, file_content, attachments_folder_id, adls_account_name, adls_container_name):
    credential = DefaultAzureCredential()
    blob_service_client = BlobServiceClient(account_url=f"https://{adls_account_name}.blob.core.windows.net", credential=credential)
    container_client = blob_service_client.get_container_client(adls_container_name)
    blob_client = container_client.get_blob_client(f"{attachments_folder_id}/{file_name}")
    try:
        blob_client.upload_blob(file_content, overwrite=True)
        logging.info(f"File '{file_name}' uploaded successfully to ADLS")
    except Exception as e:
        logging.error(f"Failed to upload {file_name} to ADLS: {e}")

def download_from_adls(attachments_folder_id, adls_account_name, adls_container_name):
    credential = DefaultAzureCredential()
    blob_service_client = BlobServiceClient(account_url=f"https://{adls_account_name}.blob.core.windows.net", credential=credential)
    container_client = blob_service_client.get_container_client(adls_container_name)
    blob_list = container_client.list_blobs(name_starts_with=f"{attachments_folder_id}/")
    file_contents = []
    for blob in blob_list:
        blob_client = container_client.get_blob_client(blob.name)
        stream = blob_client.download_blob()
        file_content = stream.readall()
        file_contents.append(process_file(blob.name, file_content))
    return file_contents

def upload_to_onedrive(file_stream, file_name, folder_id, drive_id, access_token):
    headers = {
        'Authorization': f'Bearer {access_token}',
        'Content-Type': 'application/octet-stream'
    }
    upload_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{folder_id}:/{file_name}:/content"
    response = requests.put(upload_url, headers=headers, data=file_stream)

    if response.status_code in [200, 201]:
        logging.info(f"File '{file_name}' uploaded successfully to OneDrive")
    else:
        logging.error(f"Failed to upload file '{file_name}' to OneDrive: {response.status_code}")
        logging.error(response.json())

def convert_doc_to_docx(doc_content):
    with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as temp_doc:
        temp_doc.write(doc_content)
        temp_doc_path = temp_doc.name

    temp_docx_path = temp_doc_path[:-4] + '.docx'
    word = win32.Dispatch('Word.Application')
    doc = word.Documents.Open(temp_doc_path)
    doc.SaveAs(temp_docx_path, FileFormat=12)  # 12 is the file format for .docx
    doc.Close()
    word.Quit()

    with open(temp_docx_path, 'rb') as temp_docx:
        docx_content = temp_docx.read()

    os.remove(temp_doc_path)
    os.remove(temp_docx_path)
    return docx_content
